from __future__ import annotations

import io
from typing import Any, Dict, List

from pptx import Presentation
from PIL import Image
from llama_index.core import Document

from states.loaders.utils import analyze_image_with_lvm, save_pil_image


def load_pptx(path: str, filename: str, artifacts: Dict[str, List[Any]] | None = None) -> List[Document]:
    try:
        prs = Presentation(path)
        slides_text = []
        docs: List[Document] = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                try:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                except Exception:
                    continue
            slides_text.append("\n".join(slide_text))

            for shape in slide.shapes:
                try:
                    if getattr(shape, "shape_type", None) == 13:  # Picture
                        img = shape.image
                        pil_img = Image.open(io.BytesIO(img.blob)).convert("RGB")
                        img_path = save_pil_image(pil_img, f"{filename}_slide{i + 1}")
                        caption, insights = analyze_image_with_lvm(pil_img)
                        if artifacts is not None:
                            artifacts["extracted_images"].append(img_path)
                            artifacts["image_descriptions"].append(caption)
                            artifacts["image_insights"].append(insights)
                        docs.append(
                            Document(
                                text=(
                                    "[PPTX IMAGE]\n"
                                    f"Filename:{filename}\nSlide:{i + 1}\n"
                                    f"Caption:{caption}\nInsights:{insights}"
                                ),
                                metadata={
                                    "filename": filename,
                                    "type": "pptx-image",
                                    "slide": i + 1,
                                    "image_path": img_path,
                                },
                            )
                        )
                except Exception:
                    continue

        docs.insert(
            0,
            Document(
                text="\n---SLIDES---\n".join(filter(None, slides_text)) or "[EMPTY PPTX]",
                metadata={"filename": filename, "type": "pptx"},
            ),
        )
        return docs
    except Exception as e:
        print(f"Failed to load PPTX {filename}: {e}")
        return []
